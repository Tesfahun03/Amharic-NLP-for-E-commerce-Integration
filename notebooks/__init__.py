from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1: Cover Page
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Leveraging Generative AI in Fintech to Empower 5M Young Africans in 5 Years"
subtitle.text = "A Vision for Financial Inclusion and Youth Empowerment\nPresented by: [Your Name]"

# Slide 2: The Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Addressing Youth Financial Exclusion in Africa"
content.text = (
    "Scope of the Problem:\n"
    "- 60% of Africa’s population is under 25 years old (~720M people).\n"
    "- Financial exclusion affects over 350M youth due to lack of credit access, limited financial literacy, and high unemployment (World Bank, 2022).\n\n"
    "Regional Focus:\n"
    "- Predominantly Sub-Saharan Africa, where youth unemployment rates exceed 15% (ILO, 2023).\n"
    "- Similar patterns in other developing regions but particularly acute in Africa due to infrastructure gaps.\n\n"
    "Impact:\n"
    "- Stifles entrepreneurial efforts and exacerbates poverty cycles.\n"
    "- Significant untapped potential for economic growth and innovation."
)

# Slide 3: Proposed Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "AI-Powered Financial Literacy and Credit Access Platform"
content.text = (
    "Key Features:\n"
    "- Generative AI:\n"
    "  - Personalized financial coaching via AI-driven chatbots in multiple African languages.\n"
    "  - Dynamic learning modules tailored to user needs and literacy levels.\n"
    "- Credit Scoring:\n"
    "  - AI models leveraging non-traditional data (e.g., mobile usage, transaction history) to create fairer credit scores.\n"
    "- Gamification:\n"
    "  - Interactive simulations to teach budgeting, saving, and investment in an engaging manner.\n\n"
    "Economic Feasibility:\n"
    "- Requires minimal infrastructure beyond mobile devices and internet.\n"
    "- Capitalizes on high mobile penetration rates (~46% in Sub-Saharan Africa)."
)

# Slide 4: Role of Technology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Generative AI’s Role in the Solution"
content.text = (
    "Implementation:\n"
    "- AI Models: Train models on regional financial data to ensure contextual accuracy.\n"
    "- Data Security: Utilize encryption and decentralized storage to comply with regulations.\n"
    "- Sustainability:\n"
    "  - Lightweight mobile-first applications to ensure accessibility.\n"
    "  - Integration with existing mobile money platforms (e.g., M-Pesa).\n\n"
    "Comparison to Current Solutions:\n"
    "- Traditional financial education programs are static and generic.\n"
    "- Generative AI provides personalized, adaptive, and scalable solutions."
)

# Slide 5: Improvements on Past Solutions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Innovation Through Generative AI"
content.text = (
    "Previous Approaches:\n"
    "- NGOs and government-led financial literacy campaigns (limited reach and adaptability).\n"
    "- Microfinance initiatives with traditional credit scoring (excluding informal earners).\n\n"
    "Advancements:\n"
    "- AI-enabled personalization for higher engagement.\n"
    "- Broader inclusivity through non-traditional credit assessment.\n"
    "- Scalable, multilingual platforms to reach diverse demographics."
)

# Slide 6: Funding and Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Funding Strategy and Next Steps"
content.text = (
    "Funding Sources:\n"
    "1. Development Organizations: Partnerships with entities like the World Bank or African Development Bank.\n"
    "2. Private Sector: Collaborations with fintech leaders and telecom providers.\n"
    "3. Impact Investors: Organizations focusing on social entrepreneurship and youth empowerment.\n\n"
    "Next Steps:\n"
    "- Secure initial funding for pilot projects in three countries (e.g., Kenya, Nigeria, South Africa).\n"
    "- Develop partnerships with local financial institutions and mobile operators.\n"
    "- Launch a pilot within 12 months and scale to impact 5M youth by 2030.\n\n"
    "Join Us:\n"
    "Together, we can create a financially inclusive future for Africa’s youth!"
)

# Save the presentation
prs.save("GenerativeAI_Fintech_Africa.pptx")